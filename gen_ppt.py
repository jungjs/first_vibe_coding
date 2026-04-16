import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── 색상 팔레트 ───────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
BLUE    = RGBColor(0x2E, 0x86, 0xAB)
LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY   = RGBColor(0x88, 0x88, 0x88)
BLACK   = RGBColor(0x22, 0x22, 0x22)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
RED     = RGBColor(0xC0, 0x39, 0x2B)
YELLOW  = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃

# ─── 공통 헬퍼 ────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 size=10, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, bg=None, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = '맑은 고딕'
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    return txb

def header_bar(slide, title, subtitle=''):
    """슬라이드 상단 헤더"""
    add_rect(slide, 0, 0, 13.33, 0.9, fill=NAVY)
    add_text_box(slide, title, 0.2, 0.1, 10, 0.45,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.2, 0.52, 10, 0.3,
                     size=9, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)
    # 우상단 로고 텍스트
    add_text_box(slide, 'SNTL 통합물류플랫폼', 10.5, 0.1, 2.6, 0.35,
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

def footer_bar(slide, scr_id='', page=''):
    """슬라이드 하단 푸터"""
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=NAVY)
    add_text_box(slide, 'SNTL 통합 물류 플랫폼 — UI 화면설계서  |  Confidential', 0.2, 7.15, 9, 0.3,
                 size=7, color=RGBColor(0xAA, 0xBB, 0xCC))
    if scr_id:
        add_text_box(slide, scr_id, 10, 7.15, 1.5, 0.3,
                     size=7, color=YELLOW, align=PP_ALIGN.RIGHT)
    if page:
        add_text_box(slide, page, 11.8, 7.15, 1.3, 0.3,
                     size=7, color=WHITE, align=PP_ALIGN.RIGHT)

def section_badge(slide, text, l, t, color=BLUE):
    add_rect(slide, l, t, 1.6, 0.28, fill=color)
    add_text_box(slide, text, l, t, 1.6, 0.28,
                 size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def info_box(slide, label, value, l, t, w=2.8, label_w=1.0):
    """라벨-값 인포박스"""
    add_rect(slide, l, t, label_w, 0.28, fill=BLUE)
    add_text_box(slide, label, l, t, label_w, 0.28,
                 size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, l+label_w, t, w-label_w, 0.28, fill=LIGHT, line=GRAY)
    add_text_box(slide, value, l+label_w+0.05, t, w-label_w-0.1, 0.28,
                 size=8, color=BLACK)

def wireframe_input(slide, label, l, t, w=3.5):
    """입력 필드 와이어프레임"""
    add_text_box(slide, label, l, t, w, 0.2, size=7, color=DGRAY)
    add_rect(slide, l, t+0.2, w, 0.28, fill=WHITE, line=GRAY, line_w=Pt(0.75))
    add_text_box(slide, f'{label} 입력', l+0.08, t+0.2, w-0.1, 0.28, size=7, color=GRAY)

def wireframe_btn(slide, label, l, t, w=1.5, h=0.32, color=BLUE):
    add_rect(slide, l, t, w, h, fill=color)
    add_text_box(slide, label, l, t, w, h,
                 size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def wireframe_table_header(slide, cols, l, t, total_w):
    col_w = total_w / len(cols)
    for i, col in enumerate(cols):
        add_rect(slide, l + i*col_w, t, col_w, 0.28, fill=NAVY, line=WHITE, line_w=Pt(0.5))
        add_text_box(slide, col, l + i*col_w, t, col_w, 0.28,
                     size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def wireframe_table_row(slide, cols, values, l, t, total_w, even=True):
    col_w = total_w / len(cols)
    bg = RGBColor(0xF8, 0xF9, 0xFA) if even else WHITE
    for i, val in enumerate(values):
        add_rect(slide, l + i*col_w, t, col_w, 0.24, fill=bg, line=GRAY, line_w=Pt(0.5))
        add_text_box(slide, val, l + i*col_w + 0.05, t, col_w-0.1, 0.24, size=6.5, color=BLACK)

def note_box(slide, text, l, t, w, h):
    add_rect(slide, l, t, w, h, fill=RGBColor(0xFF, 0xFF, 0xDD), line=YELLOW)
    add_text_box(slide, '📌 ' + text, l+0.1, t+0.05, w-0.2, h-0.1, size=7.5, color=RGBColor(0x66, 0x44, 0x00))

def tag(slide, text, l, t, color=GREEN):
    add_rect(slide, l, t, len(text)*0.1+0.3, 0.22, fill=color)
    add_text_box(slide, text, l+0.05, t, len(text)*0.1+0.2, 0.22, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── 슬라이드 제작 ─────────────────────────────────────────────

page = 0

def new_page():
    global page
    page += 1
    return page

# ===================================================================
# 표지
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
new_page()
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(slide, 0, 3.2, 13.33, 1.8, fill=BLUE)
add_text_box(slide, 'SNTL 통합 물류 플랫폼', 1, 1.2, 11, 1.0,
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 'UI 화면설계서', 1, 2.2, 11, 0.7,
             size=24, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text_box(slide, 'Wireframe & Screen Specification', 1, 3.35, 11, 0.5,
             size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 'AIR  ·  SEA  ·  CIR  ·  CCL', 1, 3.9, 11, 0.5,
             size=12, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_text_box(slide, 'Version 1.0  |  2026-04-16  |  SNTL', 1, 5.5, 11, 0.5,
             size=10, color=DGRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 'Confidential — Internal Use Only', 1, 6.2, 11, 0.4,
             size=9, color=DGRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# 목차
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, '목차', 'Table of Contents')
footer_bar(slide, page=f'{page}')

toc_items = [
    ('1', '로그인 / 인증', 'SCR-001 ~ SCR-005', '회원가입, 로그인, ID/PW 찾기'),
    ('2', '회원관리', 'SCR-010 ~ SCR-014', '마이페이지, 회원정보 수정, 탈퇴, 법인관리'),
    ('3', '오더관리', 'SCR-020 ~ SCR-024', '오더 목록/등록/수정, 서비스 추가, 운송비용'),
    ('4', '마스터오더관리', 'SCR-030 ~ SCR-032', '마스터오더 목록/등록/패킹/상세'),
    ('5', '창고관리', 'SCR-040 ~ SCR-042', '입고 처리, 출고 처리, 창고 현황'),
    ('6', '운송 Tracking', 'SCR-050 ~ SCR-053', 'AIR/SEA/CIR Tracking, 통관신고'),
    ('7', '회계/청구', 'SCR-060 ~ SCR-065', '청구서, 입금, 세금계산서, 수입현황, 원가'),
    ('8', 'VOC 관리', 'SCR-070 ~ SCR-072', 'VOC 목록/등록/답변'),
    ('9', '고객지원', 'SCR-080 ~ SCR-083', 'QnA, FAQ, 공지사항'),
    ('10', '관리자', 'SCR-090 ~ SCR-099', '대시보드, 회원관리, 법인심사, 시스템관리'),
]
colors_toc = [
    RGBColor(0x2E,0x86,0xAB), RGBColor(0x27,0xAE,0x60), RGBColor(0xF3,0x9C,0x12),
    RGBColor(0xE6,0x7E,0x22), RGBColor(0x8E,0x44,0xAD), RGBColor(0x16,0xA0,0x85),
    RGBColor(0xD3,0x54,0x00), RGBColor(0x2C,0x3E,0x50), RGBColor(0x7F,0x8C,0x8D),
    RGBColor(0xC0,0x39,0x2B),
]

for i, (num, title, scr_range, desc) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.5
    ty = 1.0 + row * 1.15
    c = colors_toc[i]
    add_rect(slide, lx, ty, 6.1, 0.95, fill=RGBColor(0xF5,0xF8,0xFD), line=c, line_w=Pt(1.5))
    add_rect(slide, lx, ty, 0.45, 0.95, fill=c)
    add_text_box(slide, num, lx, ty+0.25, 0.45, 0.45,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, lx+0.55, ty+0.05, 3.5, 0.38,
                 size=12, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_text_box(slide, scr_range, lx+0.55, ty+0.45, 1.8, 0.28,
                 size=8, color=c, align=PP_ALIGN.LEFT)
    add_text_box(slide, desc, lx+2.5, ty+0.47, 3.5, 0.28,
                 size=7.5, color=DGRAY, align=PP_ALIGN.LEFT)

add_text_box(slide, f'총 48개 화면', 0.4, 6.75, 12, 0.3,
             size=9, color=DGRAY, align=PP_ALIGN.RIGHT)


# ===================================================================
# 헬퍼: 섹션 타이틀 슬라이드
# ===================================================================
def section_title_slide(num, title, subtitle, color):
    sl = prs.slides.add_slide(blank_layout)
    new_page()
    add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    add_rect(sl, 0, 2.8, 13.33, 2.0, fill=color)
    add_text_box(sl, f'Section {num}', 1, 1.8, 11, 0.7,
                 size=20, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)
    add_text_box(sl, title, 1, 2.9, 11, 0.9,
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, subtitle, 1, 3.85, 11, 0.5,
                 size=12, color=YELLOW, align=PP_ALIGN.CENTER)
    footer_bar(sl, page=f'{page}')


# ===================================================================
# SECTION 1: 로그인/인증
# ===================================================================
section_title_slide('01', '로그인 / 인증', 'Login & Authentication  |  SCR-001 ~ SCR-005', BLUE)

# SCR-001 로그인
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-001  로그인', '로그인/인증 > 로그인  |  사용자: 전체')
footer_bar(slide, 'SCR-001', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))

# 메타정보
info_box(slide, '화면ID', 'SCR-001', 0.3, 1.1)
info_box(slide, '사용자', '전체', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)
info_box(slide, '상태', '미개발', 8.5, 1.1)

# 로그인 폼 (가운데 배치)
add_rect(slide, 4.0, 1.6, 5.2, 4.8, fill=WHITE, line=GRAY)
add_text_box(slide, 'SNTL', 4.0, 1.7, 5.2, 0.6, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text_box(slide, '통합 물류 플랫폼', 4.0, 2.25, 5.2, 0.3, size=9, color=DGRAY, align=PP_ALIGN.CENTER)
add_rect(slide, 4.2, 2.65, 4.8, 0.02, fill=GRAY)
wireframe_input(slide, '아이디 (ID)', 4.3, 2.75, 4.6)
wireframe_input(slide, '비밀번호 (Password)', 4.3, 3.3, 4.6)
wireframe_btn(slide, '로그인', 4.3, 3.9, 4.6, 0.38)
add_text_box(slide, 'ID 찾기  |  PW 재설정  |  회원가입', 4.0, 4.38, 5.2, 0.28, size=7.5, color=BLUE, align=PP_ALIGN.CENTER)
add_rect(slide, 4.2, 4.72, 4.8, 0.02, fill=GRAY)
add_text_box(slide, '소셜 로그인', 4.0, 4.78, 5.2, 0.25, size=7.5, color=DGRAY, align=PP_ALIGN.CENTER)
wireframe_btn(slide, 'G  Google', 4.3, 5.1, 2.1, 0.3, color=RGBColor(0xDB,0x44,0x37))
wireframe_btn(slide, 'K  Kakao', 6.5, 5.1, 2.1, 0.3, color=RGBColor(0xFE,0xE5,0x00))

# 기능 설명
add_text_box(slide, '주요 기능', 0.3, 1.5, 3.5, 0.28, size=9, bold=True, color=NAVY)
features = [
    '✓ ID/PW 입력 로그인',
    '✓ Google·Kakao OAuth2 소셜 로그인',
    '✓ 로그인 성공 시 JWT Access Token(30분) / Refresh Token(7일) 발급',
    '✓ 잘못된 자격증명 → 에러 메시지 표시',
    '✓ 비밀번호 마스킹 (표시/숨기기 토글)',
]
for j, f in enumerate(features):
    add_text_box(slide, f, 0.35, 1.85+j*0.32, 3.6, 0.28, size=7.5, color=BLACK)

note_box(slide, '입고 이후 오더 수정·삭제 불가. 창고 입고 시 오더 상태 변경됨.', 0.3, 5.65, 3.5, 0.45)
add_text_box(slide, '연결 화면: SCR-002(ID찾기), SCR-003(PW재설정), SCR-004(개인회원가입), SCR-005(법인회원가입)',
             0.3, 6.2, 12.5, 0.25, size=7, color=DGRAY)

# SCR-002 ID 찾기
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-002  ID 찾기', '로그인/인증 > ID/PW 찾기  |  사용자: 회원')
footer_bar(slide, 'SCR-002', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-002', 0.3, 1.1)
info_box(slide, '사용자', '회원', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

# 탭
add_rect(slide, 4.0, 1.55, 2.6, 0.32, fill=BLUE)
add_text_box(slide, '개인회원 ID 찾기', 4.0, 1.55, 2.6, 0.32, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 6.6, 1.55, 2.6, 0.32, fill=LIGHT, line=GRAY)
add_text_box(slide, '법인회원 ID 찾기', 6.6, 1.55, 2.6, 0.32, size=8, color=DGRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 4.0, 1.87, 5.2, 3.8, fill=WHITE, line=GRAY)
wireframe_input(slide, '이름', 4.3, 2.0, 4.6)
wireframe_input(slide, '생년월일 (YYYYMMDD)', 4.3, 2.55, 4.6)
wireframe_input(slide, '휴대폰번호', 4.3, 3.1, 3.3)
wireframe_btn(slide, '인증번호 발송', 7.65, 3.1, 1.25, 0.28, color=GREEN)
wireframe_input(slide, '인증번호 6자리', 4.3, 3.65, 4.6)
wireframe_btn(slide, 'ID 찾기', 4.3, 4.22, 4.6, 0.36)
add_text_box(slide, '* SMS 인증 완료 후 아이디 표시', 4.3, 4.65, 4.6, 0.25, size=7, color=DGRAY)

add_text_box(slide, '주요 기능', 0.3, 1.5, 3.5, 0.28, size=9, bold=True, color=NAVY)
for j, f in enumerate(['✓ 개인: 이름+생년월일+휴대폰 SMS 인증', '✓ 법인: 법인명+사업자번호 확인', '✓ 중복 계정 존재 시 복수 ID 표시']):
    add_text_box(slide, f, 0.35, 1.85+j*0.32, 3.6, 0.28, size=7.5, color=BLACK)

# SCR-004 개인회원 가입
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-004  개인회원 가입', '로그인/인증 > 회원가입  |  사용자: 개인회원')
footer_bar(slide, 'SCR-004', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-004', 0.3, 1.1)
info_box(slide, '사용자', '개인회원', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

# 스텝 인디케이터
steps = ['STEP 1\n본인확인', 'STEP 2\n정보입력', 'STEP 3\n완료']
for i, s in enumerate(steps):
    col = BLUE if i == 0 else GRAY
    add_rect(slide, 4.5+i*1.8, 1.55, 1.5, 0.5, fill=col)
    add_text_box(slide, s, 4.5+i*1.8, 1.55, 1.5, 0.5, size=7, bold=(i==0), color=WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text_box(slide, '▶', 6.05+i*1.8, 1.7, 0.3, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 2.15, 6.2, 4.3, fill=WHITE, line=GRAY)
add_text_box(slide, '본인 확인', 3.5, 2.2, 6.2, 0.35, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
wireframe_input(slide, '이름', 3.8, 2.65, 5.6)
wireframe_input(slide, '생년월일 (YYYYMMDD)', 3.8, 3.2, 5.6)
wireframe_input(slide, '휴대폰번호', 3.8, 3.75, 3.8)
wireframe_btn(slide, '인증번호 발송', 7.65, 3.75, 1.75, 0.28, color=GREEN)
wireframe_input(slide, '인증번호 입력', 3.8, 4.3, 5.6)
wireframe_btn(slide, '인증 확인', 3.8, 4.85, 2.6, 0.35)
wireframe_btn(slide, '다음 →', 6.5, 4.85, 2.9, 0.35)
note_box(slide, '동일 휴대폰번호 중복가입 체크', 3.5, 5.3, 6.2, 0.35)

add_text_box(slide, '주요 기능', 0.3, 1.5, 3.0, 0.28, size=9, bold=True, color=NAVY)
for j, f in enumerate(['✓ STEP 1: SMS 본인인증', '✓ STEP 2: ID/PW/이메일/주소 입력', '✓ STEP 3: 가입완료 + 로그인 이동', '✓ 소셜 로그인 연동 가입 지원']):
    add_text_box(slide, f, 0.35, 1.85+j*0.32, 3.0, 0.28, size=7.5, color=BLACK)

# SCR-005 법인회원 가입
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-005  법인회원 가입', '로그인/인증 > 회원가입  |  사용자: 법인관리자')
footer_bar(slide, 'SCR-005', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-005', 0.3, 1.1)
info_box(slide, '사용자', '법인관리자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

add_rect(slide, 3.5, 1.55, 6.2, 5.0, fill=WHITE, line=GRAY)
add_text_box(slide, '법인회원 가입 신청', 3.5, 1.6, 6.2, 0.38, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
wireframe_input(slide, '법인명', 3.8, 2.1, 5.6)
wireframe_input(slide, '사업자등록번호', 3.8, 2.65, 5.6)
wireframe_input(slide, '대표자명', 3.8, 3.2, 5.6)
wireframe_input(slide, '법인 이메일', 3.8, 3.75, 5.6)
wireframe_input(slide, '법인 연락처', 3.8, 4.3, 5.6)
add_rect(slide, 3.8, 4.9, 5.6, 0.5, fill=LIGHT, line=GRAY, line_w=Pt(1))
add_text_box(slide, '📎  사업자등록증 파일 첨부  (PDF/JPG, 최대 10MB)', 3.8, 4.9, 5.6, 0.5, size=8, color=DGRAY, align=PP_ALIGN.CENTER)
wireframe_btn(slide, '가입 신청', 3.8, 5.5, 5.6, 0.38)
note_box(slide, '첨부파일 업로드 후 관리자 심사 → 승인 시 계정 활성화 + 알림 발송', 0.3, 5.6, 3.0, 0.45)


# ===================================================================
# SECTION 2: 오더관리
# ===================================================================
section_title_slide('02', '오더 관리', 'Order Management  |  SCR-020 ~ SCR-024', RGBColor(0xF3,0x9C,0x12))

# SCR-020 오더 목록
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-020  오더 목록', '오더관리 > 오더 기본관리  |  사용자: 회원 / 운영자')
footer_bar(slide, 'SCR-020', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-020', 0.3, 1.1)
info_box(slide, '사용자', '회원/운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

# 검색 영역
add_rect(slide, 0.3, 1.5, 12.7, 0.75, fill=WHITE, line=GRAY, line_w=Pt(0.75))
add_text_box(slide, '검색 조건', 0.4, 1.55, 1.2, 0.25, size=8, bold=True, color=NAVY)
wireframe_input(slide, '오더번호', 1.6, 1.52, 2.0)
wireframe_input(slide, '회원명', 3.7, 1.52, 2.0)
add_rect(slide, 5.8, 1.55, 1.8, 0.32, fill=WHITE, line=GRAY, line_w=Pt(0.75))
add_text_box(slide, '서비스 ▼', 5.9, 1.55, 1.7, 0.32, size=8, color=DGRAY)
add_rect(slide, 7.7, 1.55, 1.8, 0.32, fill=WHITE, line=GRAY, line_w=Pt(0.75))
add_text_box(slide, '상태 ▼', 7.8, 1.55, 1.7, 0.32, size=8, color=DGRAY)
wireframe_btn(slide, '🔍 검색', 9.6, 1.55, 1.3, 0.32)
wireframe_btn(slide, '+ 오더 등록', 11.0, 1.55, 1.8, 0.32, color=GREEN)

# 테이블
cols = ['오더번호', '회원명', '서비스', '운송구간', '상태', '등록일', '운송비', '관리']
wireframe_table_header(slide, cols, 0.3, 2.35, 12.7)
sample_orders = [
    ('ORD-2026-00123', '홍길동', 'AIR', 'ICN→LAX', '운송중', '2026-04-10', '₩320,000', '상세'),
    ('ORD-2026-00122', '(주)ACME물산', 'SEA', 'PUS→SHA', '창고입고', '2026-04-09', '₩85,000', '상세'),
    ('ORD-2026-00121', '김영희', 'CIR', 'ICN→NRT', '등록', '2026-04-08', '₩45,000', '상세'),
    ('ORD-2026-00120', '이철수', 'CCL', 'ICN→LAX', '통관중', '2026-04-07', '₩12,000', '상세'),
    ('ORD-2026-00119', '박민준', 'AIR', 'ICN→FRA', '완료', '2026-04-06', '₩280,000', '상세'),
]
for i, vals in enumerate(sample_orders):
    wireframe_table_row(slide, cols, list(vals), 0.3, 2.63+i*0.28, 12.7, even=(i%2==0))

# 페이징
add_text_box(slide, '총 142건  |  ◀ 이전    1  2  3  ...  15    다음 ▶', 3.5, 4.1, 6.5, 0.28, size=8, color=DGRAY, align=PP_ALIGN.CENTER)
note_box(slide, '회원은 본인 오더만 조회. 운영자/관리자는 전체 조회 가능.', 0.3, 4.5, 7.0, 0.35)

# SCR-021 오더 등록
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-021  오더 등록', '오더관리 > 오더 기본관리  |  사용자: 회원')
footer_bar(slide, 'SCR-021', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-021', 0.3, 1.1)
info_box(slide, '사용자', '회원', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

# 송하인 정보
add_rect(slide, 0.3, 1.5, 6.1, 2.2, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 1.5, 6.1, 0.3, fill=NAVY)
add_text_box(slide, '📦 송하인 (Shipper) 정보', 0.4, 1.52, 5.8, 0.27, size=8, bold=True, color=WHITE)
wireframe_input(slide, '송하인명', 0.4, 1.92, 2.7)
wireframe_input(slide, '연락처', 3.2, 1.92, 3.1)
wireframe_input(slide, '국가', 0.4, 2.47, 1.3)
wireframe_input(slide, '우편번호', 1.8, 2.47, 1.5)
wireframe_input(slide, '주소', 3.4, 2.47, 2.9)
wireframe_input(slide, '상세주소', 0.4, 3.02, 5.8)

# 수하인 정보
add_rect(slide, 6.5, 1.5, 6.4, 2.2, fill=WHITE, line=GRAY)
add_rect(slide, 6.5, 1.5, 6.4, 0.3, fill=BLUE)
add_text_box(slide, '🏠 수하인 (Consignee) 정보', 6.6, 1.52, 6.0, 0.27, size=8, bold=True, color=WHITE)
wireframe_input(slide, '수하인명', 6.6, 1.92, 2.7)
wireframe_input(slide, '연락처', 9.4, 1.92, 3.4)
wireframe_input(slide, '국가', 6.6, 2.47, 1.3)
wireframe_input(slide, '우편번호', 8.0, 2.47, 1.5)
wireframe_input(slide, '주소', 9.6, 2.47, 3.2)
wireframe_input(slide, '상세주소', 6.6, 3.02, 6.2)

# 화물 정보
add_rect(slide, 0.3, 3.8, 12.6, 2.0, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 3.8, 12.6, 0.3, fill=RGBColor(0xE6,0x7E,0x22))
add_text_box(slide, '📦 화물 정보', 0.4, 3.82, 12.0, 0.27, size=8, bold=True, color=WHITE)
wireframe_input(slide, '화물명', 0.4, 4.2, 2.5)
wireframe_input(slide, 'HS코드', 3.0, 4.2, 2.0)
wireframe_input(slide, '총중량 (kg)', 5.1, 4.2, 2.0)
wireframe_input(slide, 'CBM', 7.2, 4.2, 1.5)
wireframe_input(slide, '포장수량', 8.8, 4.2, 1.5)
wireframe_input(slide, '화물가액 (USD)', 10.4, 4.2, 2.4)
wireframe_input(slide, '특이사항 / 비고', 0.4, 4.75, 12.3)

wireframe_btn(slide, '저장', 9.5, 5.85, 1.5, 0.38)
wireframe_btn(slide, '서비스 추가 →', 11.1, 5.85, 2.0, 0.38, color=GREEN)
wireframe_btn(slide, '취소', 7.9, 5.85, 1.4, 0.38, color=RED)

# SCR-022 오더 상세/수정
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-022  오더 상세 / 수정', '오더관리 > 오더 기본관리  |  사용자: 회원 / 운영자')
footer_bar(slide, 'SCR-022', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-022', 0.3, 1.1)
info_box(slide, '사용자', '회원/운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

# 오더 헤더
add_rect(slide, 0.3, 1.5, 12.7, 0.5, fill=WHITE, line=GRAY)
add_text_box(slide, '오더번호: ORD-2026-00123', 0.5, 1.55, 4.0, 0.38, size=11, bold=True, color=NAVY)
tag(slide, '운송중', 4.6, 1.62, color=BLUE)
tag(slide, 'AIR', 5.4, 1.62, color=RGBColor(0x8E,0x44,0xAD))
add_text_box(slide, '등록일: 2026-04-10 14:32', 8.0, 1.6, 3.5, 0.28, size=8, color=DGRAY, align=PP_ALIGN.RIGHT)

# 3 컬럼 정보
for col_i, (title_, fields) in enumerate([
    ('송하인 정보', ['이름: 홍길동', '연락처: 010-1234-5678', '국가: KR', '주소: 서울 강남구 테헤란로 123']),
    ('수하인 정보', ['이름: John Smith', '연락처: +1-310-xxx-xxxx', '국가: US', '주소: 123 Main St, Los Angeles']),
    ('화물 정보', ['화물명: 전자부품', 'HS코드: 8542.31', '중량: 25.5 kg', 'CBM: 0.12']),
]):
    lx = 0.3 + col_i * 4.3
    add_rect(slide, lx, 2.1, 4.1, 2.0, fill=WHITE, line=GRAY)
    add_rect(slide, lx, 2.1, 4.1, 0.28, fill=NAVY if col_i==0 else (BLUE if col_i==1 else RGBColor(0xE6,0x7E,0x22)))
    add_text_box(slide, title_, lx+0.1, 2.12, 3.9, 0.25, size=8, bold=True, color=WHITE)
    for j, field in enumerate(fields):
        add_text_box(slide, field, lx+0.15, 2.45+j*0.35, 3.9, 0.3, size=8, color=BLACK)

# 서비스 정보 탭
add_rect(slide, 0.3, 4.2, 12.7, 1.9, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 4.2, 12.7, 0.28, fill=BLUE)
add_text_box(slide, '서비스 정보', 0.5, 4.22, 5.0, 0.25, size=8, bold=True, color=WHITE)
cols_svc = ['서비스유형', '출발지', '목적지', '출발일', '도착예정일', '운송비', '상태']
wireframe_table_header(slide, cols_svc, 0.3, 4.5, 12.7)
wireframe_table_row(slide, cols_svc, ['AIR', 'ICN', 'LAX', '2026-04-12', '2026-04-14', '₩320,000', '운송중'], 0.3, 4.78, 12.7)

wireframe_btn(slide, '수정', 9.5, 6.18, 1.4, 0.35, color=ORANGE)
wireframe_btn(slide, '삭제', 11.0, 6.18, 1.2, 0.35, color=RED)
wireframe_btn(slide, '+ 서비스 추가', 7.0, 6.18, 2.3, 0.35, color=GREEN)
note_box(slide, '입고(창고 수령) 이후 수정·삭제 불가. 버튼 비활성화 처리.', 0.3, 6.18, 6.5, 0.38)

# SCR-023 서비스 추가
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-023  서비스 추가', '오더관리 > 서비스 관리  |  사용자: 회원')
footer_bar(slide, 'SCR-023', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-023', 0.3, 1.1)
info_box(slide, '사용자', '회원', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

# 서비스 선택 탭
svc_tabs = [('AIR 항운', NAVY, True), ('SEA 해운', GRAY, False), ('CIR 택배', GRAY, False), ('CCL 통관', GRAY, False)]
for i, (t, c, act) in enumerate(svc_tabs):
    add_rect(slide, 2.5+i*2.2, 1.55, 2.1, 0.38, fill=c if not act else BLUE)
    add_text_box(slide, t, 2.5+i*2.2, 1.55, 2.1, 0.38, size=9, bold=act, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 2.5, 1.93, 8.8, 4.5, fill=WHITE, line=GRAY)
add_text_box(slide, 'AIR 항운 서비스 정보', 2.5, 1.98, 8.8, 0.35, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
wireframe_input(slide, '출발 공항 (IATA코드)', 2.8, 2.45, 3.8)
wireframe_input(slide, '도착 공항 (IATA코드)', 6.8, 2.45, 4.3)
wireframe_input(slide, '항공사', 2.8, 3.0, 3.8)
wireframe_input(slide, 'MAWB번호 (Master Air Waybill)', 6.8, 3.0, 4.3)
wireframe_input(slide, 'HAWB번호 (House Air Waybill)', 2.8, 3.55, 3.8)
wireframe_input(slide, '출발 예정일', 6.8, 3.55, 4.3)
wireframe_input(slide, '도착 예정일', 2.8, 4.1, 3.8)
wireframe_input(slide, '특이사항', 6.8, 4.1, 4.3)
wireframe_btn(slide, '저장', 9.0, 5.5, 2.1, 0.38)
wireframe_btn(slide, '취소', 6.8, 5.5, 2.0, 0.38, color=RED)
note_box(slide, '서비스 유형별로 입력 항목이 다름. 하나의 오더에 최대 4가지 서비스 추가 가능.', 0.3, 5.65, 2.0, 0.5)


# ===================================================================
# SECTION 3: 마스터오더 / 창고관리
# ===================================================================
section_title_slide('03', '마스터오더 / 창고관리', 'Master Order & Warehouse  |  SCR-030 ~ SCR-042', RGBColor(0x8E,0x44,0xAD))

# SCR-031 마스터오더 등록/패킹
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-031  마스터오더 등록 / 패킹', '마스터오더관리 > 마스터오더 등록  |  사용자: 운영자')
footer_bar(slide, 'SCR-031', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-031', 0.3, 1.1)
info_box(slide, '사용자', '운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

# 왼쪽: 마스터오더 정보
add_rect(slide, 0.3, 1.5, 5.5, 5.5, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 1.5, 5.5, 0.3, fill=NAVY)
add_text_box(slide, '마스터오더 기본정보', 0.4, 1.52, 5.3, 0.27, size=8, bold=True, color=WHITE)
wireframe_input(slide, '마스터오더번호 (자동생성)', 0.4, 1.9, 5.2)
wireframe_input(slide, '운송 서비스 유형', 0.4, 2.45, 5.2)
wireframe_input(slide, '출발지', 0.4, 3.0, 2.3)
wireframe_input(slide, '목적지', 2.9, 3.0, 2.7)
wireframe_input(slide, '출발 예정일', 0.4, 3.55, 2.3)
wireframe_input(slide, '도착 예정일', 2.9, 3.55, 2.7)
wireframe_input(slide, '비고', 0.4, 4.1, 5.2)
wireframe_btn(slide, '마스터오더 생성', 0.4, 4.7, 5.2, 0.38, color=GREEN)

# 오른쪽: 오더 패킹
add_rect(slide, 5.95, 1.5, 7.05, 5.5, fill=WHITE, line=GRAY)
add_rect(slide, 5.95, 1.5, 7.05, 0.3, fill=BLUE)
add_text_box(slide, '오더 패킹 (Order Packing)', 6.05, 1.52, 6.8, 0.27, size=8, bold=True, color=WHITE)
add_text_box(slide, '오더 검색', 6.05, 1.88, 1.2, 0.25, size=7, bold=True, color=NAVY)
wireframe_input(slide, '오더번호 또는 회원명', 6.05, 2.08, 4.5)
wireframe_btn(slide, '검색', 10.65, 2.08, 1.2, 0.28, color=BLUE)
cols_pack = ['선택', '오더번호', '회원명', '서비스', '중량']
wireframe_table_header(slide, cols_pack, 6.05, 2.44, 6.8)
pack_rows = [
    ('☐', 'ORD-0124', '홍길동', 'AIR', '12kg'),
    ('☐', 'ORD-0125', '(주)ACME', 'AIR', '45kg'),
    ('☑', 'ORD-0126', '김영희', 'AIR', '8kg'),
    ('☑', 'ORD-0127', '이철수', 'AIR', '22kg'),
]
for i, row_vals in enumerate(pack_rows):
    wireframe_table_row(slide, cols_pack, list(row_vals), 6.05, 2.72+i*0.26, 6.8, even=(i%2==0))
wireframe_btn(slide, '+ 패킹 추가', 6.05, 3.82, 3.2, 0.32, color=GREEN)
wireframe_btn(slide, '패킹 제거', 9.35, 3.82, 2.6, 0.32, color=RED)
add_text_box(slide, '현재 패킹된 오더: 2개  /  총 중량: 30kg', 6.05, 4.22, 6.8, 0.28, size=7.5, color=NAVY, align=PP_ALIGN.CENTER)
note_box(slide, '패킹된 오더는 단독 수정·삭제 불가. 마스터오더 단위로 처리.', 6.05, 4.6, 6.8, 0.38)
wireframe_btn(slide, '패킹 완료 저장', 8.0, 5.1, 3.8, 0.38, color=NAVY)

# SCR-040 입고 처리
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-040  창고 입고 처리', '창고관리 > 입고 관리  |  사용자: 운영자')
footer_bar(slide, 'SCR-040', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-040', 0.3, 1.1)
info_box(slide, '사용자', '운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

# 바코드 스캔 영역
add_rect(slide, 0.3, 1.5, 6.0, 5.4, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 1.5, 6.0, 0.3, fill=NAVY)
add_text_box(slide, '바코드 스캔 / 오더 직접 입력', 0.4, 1.52, 5.7, 0.27, size=8, bold=True, color=WHITE)
add_rect(slide, 0.8, 1.95, 5.0, 2.5, fill=LIGHT, line=BLUE, line_w=Pt(2))
add_text_box(slide, '[ 바코드 스캔 영역 ]', 0.8, 2.8, 5.0, 0.8, size=14, color=DGRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, '📷 카메라/스캐너 연동', 0.8, 3.2, 5.0, 0.5, size=9, color=DGRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, '— 또는 직접 입력 —', 0.8, 4.52, 5.0, 0.28, size=8, color=DGRAY, align=PP_ALIGN.CENTER)
wireframe_input(slide, '오더번호 직접 입력', 0.8, 4.82, 3.8)
wireframe_btn(slide, '조회', 4.65, 4.82, 1.3, 0.28, color=BLUE)

# 입고 확인 영역
add_rect(slide, 6.5, 1.5, 6.4, 5.4, fill=WHITE, line=GRAY)
add_rect(slide, 6.5, 1.5, 6.4, 0.3, fill=GREEN)
add_text_box(slide, '입고 확인 정보', 6.6, 1.52, 6.2, 0.27, size=8, bold=True, color=WHITE)
info_box(slide, '오더번호', 'ORD-2026-00123', 6.6, 1.9, 6.2, 2.2)
info_box(slide, '마스터오더', 'MORD-2026-001', 6.6, 2.22, 6.2, 2.2)
info_box(slide, '회원명', '홍길동', 6.6, 2.54, 6.2, 2.2)
info_box(slide, '서비스', 'AIR (ICN→LAX)', 6.6, 2.86, 6.2, 2.2)
info_box(slide, '화물명', '전자부품 12kg', 6.6, 3.18, 6.2, 2.2)
info_box(slide, '현재 상태', '등록', 6.6, 3.5, 6.2, 2.2)
wireframe_input(slide, '입고 담당자 메모 (선택)', 6.6, 3.92, 6.2)
wireframe_btn(slide, '✓  입고 처리 확인', 6.6, 4.62, 6.2, 0.45, color=GREEN)
note_box(slide, '입고 처리 후 오더 상태 → "창고입고". 이후 수정·삭제 불가.', 6.6, 5.15, 6.2, 0.4)

# SCR-041 출고 처리
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-041  창고 출고 처리', '창고관리 > 출고 관리  |  사용자: 운영자')
footer_bar(slide, 'SCR-041', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-041', 0.3, 1.1)
info_box(slide, '사용자', '운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

add_rect(slide, 0.3, 1.5, 6.0, 5.4, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 1.5, 6.0, 0.3, fill=NAVY)
add_text_box(slide, '출고 대상 목록', 0.4, 1.52, 5.7, 0.27, size=8, bold=True, color=WHITE)
cols_wh = ['선택', '마스터오더', '오더수', '서비스', '운송장출력']
wireframe_table_header(slide, cols_wh, 0.3, 1.9, 6.0)
for i, r_ in enumerate([('☑','MORD-001','3','AIR','출력전'),('☐','MORD-002','5','SEA','출력전')]):
    wireframe_table_row(slide, cols_wh, list(r_), 0.3, 2.18+i*0.26, 6.0, even=(i%2==0))
wireframe_btn(slide, '🖨 운송장 출력', 0.4, 2.82, 2.8, 0.35, color=BLUE)
wireframe_btn(slide, '출고 처리', 3.3, 2.82, 2.8, 0.35, color=GREEN)

add_rect(slide, 6.5, 1.5, 6.4, 5.4, fill=WHITE, line=GRAY)
add_rect(slide, 6.5, 1.5, 6.4, 0.3, fill=ORANGE)
add_text_box(slide, '운송장 미리보기', 6.6, 1.52, 6.2, 0.27, size=8, bold=True, color=WHITE)
add_rect(slide, 6.7, 1.9, 6.0, 4.3, fill=WHITE, line=GRAY, line_w=Pt(1.5))
add_text_box(slide, 'SNTL LOGISTICS', 6.7, 2.0, 6.0, 0.4, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text_box(slide, 'AIR WAYBILL', 6.7, 2.4, 6.0, 0.3, size=9, color=DGRAY, align=PP_ALIGN.CENTER)
for i, line in enumerate(['발송인: 홍길동  KR', 'AWBN: 123-4567-8901', '수취인: John Smith  US', '내용물: 전자부품  12kg', '서비스: AIR EXPRESS']):
    add_text_box(slide, line, 6.8, 2.8+i*0.38, 5.8, 0.3, size=8, color=BLACK)
add_rect(slide, 7.5, 4.8, 4.5, 0.7, fill=BLACK)
add_text_box(slide, '||||||||  ORD-2026-00123  ||||||||', 7.5, 4.85, 4.5, 0.6, size=7, color=WHITE, align=PP_ALIGN.CENTER)


# ===================================================================
# SECTION 4: Tracking / 회계
# ===================================================================
section_title_slide('04', '운송 Tracking / 회계·청구', 'Tracking & Billing  |  SCR-050 ~ SCR-065', RGBColor(0x16,0xA0,0x85))

# SCR-050 AIR Tracking
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-050  AIR Tracking', '운송Tracking > 항운(AIR) Tracking  |  사용자: 회원 / 운영자')
footer_bar(slide, 'SCR-050', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-050', 0.3, 1.1)
info_box(slide, '사용자', '회원/운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

add_rect(slide, 0.3, 1.5, 12.7, 0.5, fill=WHITE, line=GRAY)
wireframe_input(slide, 'AWB번호 또는 오더번호', 0.4, 1.52, 4.0)
wireframe_btn(slide, '🔍 조회', 4.5, 1.55, 1.5, 0.35, color=BLUE)
tag(slide, 'AIR', 6.2, 1.62, color=RGBColor(0x8E,0x44,0xAD))
tag(slide, 'ICN → LAX', 6.8, 1.62, color=BLUE)
tag(slide, '운송중', 8.2, 1.62, color=ORANGE)

# 타임라인
add_rect(slide, 0.3, 2.1, 12.7, 3.8, fill=WHITE, line=GRAY)
add_text_box(slide, '운송 현황 타임라인', 0.5, 2.15, 4.0, 0.3, size=9, bold=True, color=NAVY)
tl_items = [
    ('✅', '2026-04-10 09:00', '화물 접수', 'ICN 창고 입고 완료', GREEN),
    ('✅', '2026-04-11 14:30', '출항', 'KE 123편 인천공항 출발', GREEN),
    ('🔄', '2026-04-13 08:00', '경유지 도착', 'ORD 공항 트랜짓 처리 중', ORANGE),
    ('⬜', '2026-04-14 예정', '목적지 도착', 'LAX 공항 도착 예정', GRAY),
    ('⬜', '2026-04-15 예정', '통관', '세관 통관 처리 예정', GRAY),
    ('⬜', '2026-04-15 예정', '배송완료', '수취인 인도 예정', GRAY),
]
for i, (icon, dt, title_, desc_, c) in enumerate(tl_items):
    y = 2.55 + i * 0.48
    add_rect(slide, 0.5, y, 0.3, 0.3, fill=c)
    add_text_box(slide, icon, 0.5, y, 0.3, 0.3, size=9, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, dt, 0.9, y+0.02, 2.2, 0.26, size=7.5, color=DGRAY)
    add_text_box(slide, title_, 3.2, y+0.02, 2.5, 0.26, size=8, bold=True, color=BLACK)
    add_text_box(slide, desc_, 5.8, y+0.02, 7.0, 0.26, size=7.5, color=DGRAY)
    if i < 5:
        add_rect(slide, 0.63, y+0.3, 0.04, 0.18, fill=GRAY)

note_box(slide, '외부 항공사 API 연동으로 실시간 조회. 5분 주기 자동 갱신.', 0.3, 5.98, 12.7, 0.35)

# SCR-060 청구서 목록
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-060  청구서 목록', '회계/청구 > 청구서 관리  |  사용자: 운영자 / 회원')
footer_bar(slide, 'SCR-060', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-060', 0.3, 1.1)
info_box(slide, '사용자', '운영자/회원', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

add_rect(slide, 0.3, 1.5, 12.7, 0.6, fill=WHITE, line=GRAY)
wireframe_input(slide, '청구서번호', 0.4, 1.53, 2.0)
wireframe_input(slide, '회원명', 2.5, 1.53, 2.0)
wireframe_input(slide, '기간 FROM', 4.6, 1.53, 1.8)
wireframe_input(slide, '기간 TO', 6.5, 1.53, 1.8)
add_rect(slide, 8.4, 1.55, 1.8, 0.32, fill=WHITE, line=GRAY)
add_text_box(slide, '상태 ▼', 8.5, 1.55, 1.7, 0.32, size=8, color=DGRAY)
wireframe_btn(slide, '검색', 10.3, 1.55, 1.2, 0.32)
wireframe_btn(slide, '청구서 생성', 11.6, 1.55, 1.5, 0.32, color=GREEN)

cols_inv = ['청구서번호', '회원명', '오더수', '청구금액', '입금액', '미수금', '청구일', '상태', '관리']
wireframe_table_header(slide, cols_inv, 0.3, 2.2, 12.7)
inv_rows = [
    ('INV-2026-0045','홍길동','3','₩650,000','₩650,000','₩0','2026-04-01','완납','상세'),
    ('INV-2026-0044','(주)ACME','12','₩3,200,000','₩1,600,000','₩1,600,000','2026-04-01','미납','상세'),
    ('INV-2026-0043','김영희','1','₩85,000','₩0','₩85,000','2026-03-31','미납','상세'),
]
for i, row_vals in enumerate(inv_rows):
    wireframe_table_row(slide, cols_inv, list(row_vals), 0.3, 2.48+i*0.28, 12.7, even=(i%2==0))

# 요약
add_rect(slide, 0.3, 3.4, 12.7, 0.8, fill=WHITE, line=GRAY)
for i, (label, val, color_) in enumerate([
    ('총 청구금액', '₩42,850,000', NAVY),
    ('총 입금액', '₩38,200,000', GREEN),
    ('총 미수금', '₩4,650,000', RED),
    ('청구서 건수', '45건', BLUE),
]):
    lx = 0.5 + i * 3.15
    add_rect(slide, lx, 3.45, 2.9, 0.65, fill=color_, line=None)
    add_text_box(slide, label, lx, 3.5, 2.9, 0.25, size=7.5, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, val, lx, 3.75, 2.9, 0.35, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ===================================================================
# SECTION 5: 관리자
# ===================================================================
section_title_slide('05', '관리자 시스템', 'Admin System  |  SCR-090 ~ SCR-099', RGBColor(0xC0,0x39,0x2B))

# SCR-090 관리자 대시보드
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-090  관리자 대시보드', '관리자 > 통계 관리  |  사용자: 관리자')
footer_bar(slide, 'SCR-090', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-090', 0.3, 1.1)
info_box(slide, '사용자', '관리자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

# KPI 카드 4개
kpi_data = [
    ('이번달 오더', '142건', '▲ 12%', NAVY),
    ('이번달 수익', '₩32.5M', '▲ 8.3%', GREEN),
    ('처리중 오더', '37건', '─ 0%', ORANGE),
    ('미수금', '₩4.65M', '▼ 2.1%', RED),
]
for i, (label, val, trend, c) in enumerate(kpi_data):
    lx = 0.3 + i * 3.15
    add_rect(slide, lx, 1.55, 2.95, 1.1, fill=WHITE, line=c, line_w=Pt(2))
    add_rect(slide, lx, 1.55, 2.95, 0.08, fill=c)
    add_text_box(slide, label, lx+0.1, 1.65, 2.75, 0.3, size=9, color=DGRAY)
    add_text_box(slide, val, lx+0.1, 1.95, 2.75, 0.45, size=18, bold=True, color=c)
    add_text_box(slide, trend, lx+0.1, 2.42, 2.75, 0.2, size=8, color=DGRAY)

# 차트 영역 2개
add_rect(slide, 0.3, 2.8, 6.1, 4.0, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 2.8, 6.1, 0.3, fill=NAVY)
add_text_box(slide, '월별 오더/수익 현황 (Bar Chart)', 0.4, 2.82, 5.8, 0.27, size=8, bold=True, color=WHITE)
for i, (mo, h1, h2) in enumerate([('1월',1.2,0.9),('2월',1.5,1.1),('3월',1.8,1.4),('4월',2.1,1.7),('5월',1.6,1.3),('6월',1.9,1.5)]):
    bx = 0.7 + i*0.85
    add_rect(slide, bx, 3.2+(2.1-h1), 0.35, h1, fill=NAVY)
    add_rect(slide, bx+0.38, 3.2+(2.1-h2), 0.35, h2, fill=BLUE)
    add_text_box(slide, mo, bx, 5.35, 0.7, 0.22, size=6.5, color=DGRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 6.6, 2.8, 6.3, 4.0, fill=WHITE, line=GRAY)
add_rect(slide, 6.6, 2.8, 6.3, 0.3, fill=BLUE)
add_text_box(slide, '서비스별 오더 비중 (Pie Chart)', 6.7, 2.82, 6.0, 0.27, size=8, bold=True, color=WHITE)
# 도넛 차트 모사
add_rect(slide, 7.5, 3.2, 2.5, 2.5, fill=LIGHT, line=GRAY)
add_text_box(slide, '[ Pie Chart ]', 7.5, 4.0, 2.5, 0.8, size=9, color=DGRAY, align=PP_ALIGN.CENTER)
for i, (svc, pct, c_) in enumerate([('AIR','45%',NAVY),('SEA','28%',BLUE),('CIR','18%',GREEN),('CCL','9%',ORANGE)]):
    add_rect(slide, 10.2, 3.3+i*0.55, 0.25, 0.25, fill=c_)
    add_text_box(slide, f'{svc}: {pct}', 10.55, 3.3+i*0.55, 2.2, 0.28, size=8, color=BLACK)

# SCR-091 회원 목록
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-091  회원 목록 / 관리', '관리자 > 회원관리  |  사용자: 관리자')
footer_bar(slide, 'SCR-091', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-091', 0.3, 1.1)
info_box(slide, '사용자', '관리자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

add_rect(slide, 0.3, 1.5, 12.7, 0.6, fill=WHITE, line=GRAY)
wireframe_input(slide, '회원명 또는 ID', 0.4, 1.53, 2.5)
wireframe_input(slide, '이메일', 3.0, 1.53, 2.5)
add_rect(slide, 5.6, 1.55, 1.8, 0.32, fill=WHITE, line=GRAY)
add_text_box(slide, '회원유형 ▼', 5.7, 1.55, 1.7, 0.32, size=8, color=DGRAY)
add_rect(slide, 7.5, 1.55, 1.8, 0.32, fill=WHITE, line=GRAY)
add_text_box(slide, '상태 ▼', 7.6, 1.55, 1.7, 0.32, size=8, color=DGRAY)
wireframe_btn(slide, '검색', 9.4, 1.55, 1.2, 0.32)
wireframe_btn(slide, '엑셀 다운로드', 10.7, 1.55, 2.1, 0.32, color=GREEN)

cols_m = ['회원ID', '회원명', '유형', '이메일', '가입일', '등급', '상태', '관리']
wireframe_table_header(slide, cols_m, 0.3, 2.25, 12.7)
mem_rows = [
    ('MEM-0001','홍길동','개인','hong@email.com','2026-01-15','BRONZE','정상','상세'),
    ('MEM-0002','(주)ACME물산','법인','admin@acme.com','2026-02-10','GOLD','정상','상세'),
    ('MEM-0003','김영희','개인','kim@email.com','2026-03-05','SILVER','정상','상세'),
    ('MEM-0004','(주)신흥무역','법인','info@shin.com','2026-04-01','—','심사중','심사'),
    ('MEM-0005','이철수','개인','lee@email.com','2026-04-12','BRONZE','정지','상세'),
]
for i, row_vals in enumerate(mem_rows):
    wireframe_table_row(slide, cols_m, list(row_vals), 0.3, 2.53+i*0.28, 12.7, even=(i%2==0))

note_box(slide, '법인회원 상태 "심사중" → SCR-092 법인심사 화면으로 이동', 0.3, 4.05, 7.0, 0.35)

# SCR-092 법인심사
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-092  법인회원 심사', '관리자 > 회원관리  |  사용자: 관리자')
footer_bar(slide, 'SCR-092', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-092', 0.3, 1.1)
info_box(slide, '사용자', '관리자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 1', 5.8, 1.1)

add_rect(slide, 0.3, 1.5, 5.8, 5.4, fill=WHITE, line=GRAY)
add_rect(slide, 0.3, 1.5, 5.8, 0.3, fill=NAVY)
add_text_box(slide, '신청 법인 정보', 0.4, 1.52, 5.6, 0.27, size=8, bold=True, color=WHITE)
for j, (label_, val_) in enumerate([
    ('법인명','(주)신흥무역'),('사업자번호','123-45-67890'),
    ('대표자','박대표'),('법인이메일','info@shin.com'),
    ('신청일','2026-04-12'),('상태','심사중'),
]):
    info_box(slide, label_, val_, 0.4, 1.9+j*0.38, 5.6)

add_rect(slide, 6.3, 1.5, 6.6, 5.4, fill=WHITE, line=GRAY)
add_rect(slide, 6.3, 1.5, 6.6, 0.3, fill=BLUE)
add_text_box(slide, '사업자등록증 첨부파일', 6.4, 1.52, 6.4, 0.27, size=8, bold=True, color=WHITE)
add_rect(slide, 6.4, 1.9, 6.4, 3.5, fill=LIGHT, line=GRAY)
add_text_box(slide, '[ PDF 미리보기 영역 ]', 6.4, 3.0, 6.4, 0.8, size=11, color=DGRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, '사업자등록증.pdf', 6.4, 5.45, 3.5, 0.3, size=8, color=BLUE)
wireframe_btn(slide, '파일 다운로드', 9.9, 5.45, 2.8, 0.32, color=BLUE)

wireframe_btn(slide, '✅  승인', 0.4, 6.3, 2.5, 0.42, color=GREEN)
wireframe_btn(slide, '❌  거부', 3.1, 6.3, 2.5, 0.42, color=RED)
add_rect(slide, 6.3, 6.3, 6.6, 0.42, fill=WHITE, line=GRAY)
add_text_box(slide, '거부 사유 입력 (거부 시 필수)', 6.4, 6.32, 6.4, 0.38, size=8, color=DGRAY)

note_box(slide, '승인 시 법인 계정 활성화 + 담당자 이메일/앱 알림 자동 발송', 0.3, 6.78, 12.7, 0.35)


# ===================================================================
# SECTION 6: VOC / 고객지원
# ===================================================================
section_title_slide('06', 'VOC / 고객지원', 'VOC & Support  |  SCR-070 ~ SCR-083', RGBColor(0x2C,0x3E,0x50))

# SCR-070 VOC 목록
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, 'SCR-070  VOC 목록', 'VOC관리 > VOC 조회  |  사용자: 회원 / 운영자')
footer_bar(slide, 'SCR-070', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.0, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
info_box(slide, '화면ID', 'SCR-070', 0.3, 1.1)
info_box(slide, '사용자', '회원/운영자', 3.2, 1.1)
info_box(slide, 'Phase', 'Phase 2', 5.8, 1.1)

add_rect(slide, 0.3, 1.5, 12.7, 0.55, fill=WHITE, line=GRAY)
wireframe_input(slide, '오더번호', 0.4, 1.53, 2.5)
add_rect(slide, 3.0, 1.55, 1.8, 0.32, fill=WHITE, line=GRAY)
add_text_box(slide, '처리상태 ▼', 3.1, 1.55, 1.7, 0.32, size=8, color=DGRAY)
wireframe_input(slide, '등록일 FROM', 5.0, 1.53, 1.8)
wireframe_input(slide, '등록일 TO', 6.9, 1.53, 1.8)
wireframe_btn(slide, '검색', 8.8, 1.55, 1.2, 0.32)
wireframe_btn(slide, '+ VOC 등록', 10.1, 1.55, 1.5, 0.32, color=GREEN)

cols_v = ['VOC번호', '오더번호', '등록자', '제목', '등록일', '처리상태', '관리']
wireframe_table_header(slide, cols_v, 0.3, 2.18, 12.7)
voc_rows = [
    ('VOC-001','ORD-0123','홍길동','배송 지연 문의','2026-04-10','IN_PROGRESS','답변'),
    ('VOC-002','ORD-0115','김영희','화물 파손 신고','2026-04-08','OPEN','답변'),
    ('VOC-003','ORD-0109','이철수','운송비 오청구','2026-04-05','CLOSED','상세'),
]
for i, row_vals in enumerate(voc_rows):
    wireframe_table_row(slide, cols_v, list(row_vals), 0.3, 2.46+i*0.28, 12.7, even=(i%2==0))
note_box(slide, '회원은 본인 VOC만 조회 가능. 운영자는 전체 조회/처리 가능.', 0.3, 3.4, 12.7, 0.35)


# ===================================================================
# 화면 흐름도 (Flow Diagram)
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, '화면 흐름도', 'Screen Flow Diagram — 회원 주요 업무 흐름')
footer_bar(slide, '', f'{page}')

add_rect(slide, 0.2, 0.95, 12.9, 6.1, fill=LIGHT, line=GRAY, line_w=Pt(0.5))
add_text_box(slide, '회원 주요 업무 흐름 (User Flow)', 0.3, 1.0, 12.7, 0.38, size=11, bold=True, color=NAVY)

flows = [
    (0.5, 1.55, 'SCR-004\n개인회원가입', BLUE),
    (2.8, 1.55, 'SCR-001\n로그인', NAVY),
    (5.1, 1.55, 'SCR-020\n오더 목록', RGBColor(0xF3,0x9C,0x12)),
    (7.4, 1.55, 'SCR-021\n오더 등록', RGBColor(0xE6,0x7E,0x22)),
    (9.7, 1.55, 'SCR-023\n서비스 추가', GREEN),
    (9.7, 3.2, 'SCR-040\n창고 입고', RGBColor(0x8E,0x44,0xAD)),
    (7.4, 3.2, 'SCR-041\n창고 출고', RGBColor(0x16,0xA0,0x85)),
    (5.1, 3.2, 'SCR-050\nAIR Tracking', BLUE),
    (2.8, 3.2, 'SCR-060\n청구서', RGBColor(0xC0,0x39,0x2B)),
    (0.5, 3.2, 'SCR-070\nVOC', RGBColor(0x2C,0x3E,0x50)),
]

for lx, ty, txt, c in flows:
    add_rect(slide, lx, ty, 1.9, 0.9, fill=c)
    add_text_box(slide, txt, lx, ty+0.1, 1.9, 0.7, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 화살표 (텍스트로 대체)
arrows_h = [(0.5+1.9, 1.55+0.45), (2.8+1.9, 1.55+0.45), (5.1+1.9, 1.55+0.45), (7.4+1.9, 1.55+0.45)]
for ax, ay in arrows_h:
    add_text_box(slide, '→', ax, ay-0.15, 0.4, 0.3, size=12, color=DGRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, '↓', 11.0, 2.5, 0.4, 0.5, size=14, color=DGRAY, align=PP_ALIGN.CENTER)

arrows_h2 = [(7.4+1.9, 3.2+0.45), (5.1+1.9, 3.2+0.45), (2.8+1.9, 3.2+0.45), (0.5+1.9, 3.2+0.45)]
for ax, ay in arrows_h2:
    add_text_box(slide, '←', ax, ay-0.15, 0.4, 0.3, size=12, color=DGRAY, align=PP_ALIGN.CENTER)

# 관리자 흐름
add_text_box(slide, '관리자 흐름:', 0.3, 4.5, 2.2, 0.3, size=9, bold=True, color=NAVY)
admin_flows = [
    (0.5, 4.85, 'SCR-090\n대시보드', NAVY),
    (2.8, 4.85, 'SCR-091\n회원관리', RED),
    (5.1, 4.85, 'SCR-092\n법인심사', ORANGE),
    (7.4, 4.85, 'SCR-031\n마스터오더', RGBColor(0x8E,0x44,0xAD)),
    (9.7, 4.85, 'SCR-065\n원가관리', GREEN),
]
for lx, ty, txt, c in admin_flows:
    add_rect(slide, lx, ty, 1.9, 0.8, fill=c)
    add_text_box(slide, txt, lx, ty+0.05, 1.9, 0.7, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i in range(4):
    add_text_box(slide, '→', 0.5+1.9+i*2.3, 4.85+0.3, 0.4, 0.3, size=12, color=DGRAY, align=PP_ALIGN.CENTER)


# ===================================================================
# 마지막 페이지: 화면 목록 요약
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
new_page()
header_bar(slide, '화면 목록 요약', 'Screen List Summary  |  총 48개 화면')
footer_bar(slide, '', f'{page}')
add_rect(slide, 0.2, 1.0, 12.9, 6.1, fill=LIGHT, line=GRAY, line_w=Pt(0.5))

summary = [
    ('로그인/인증', ['SCR-001 로그인','SCR-002 ID찾기','SCR-003 PW재설정','SCR-004 개인회원가입','SCR-005 법인회원가입'], BLUE),
    ('회원관리', ['SCR-010 마이페이지','SCR-011 회원정보수정','SCR-012 회원탈퇴','SCR-013 법인마이페이지','SCR-014 부서/멤버관리'], GREEN),
    ('오더관리', ['SCR-020 오더목록','SCR-021 오더등록','SCR-022 오더상세/수정','SCR-023 서비스추가','SCR-024 운송비용조회'], RGBColor(0xF3,0x9C,0x12)),
    ('마스터오더', ['SCR-030 마스터오더목록','SCR-031 마스터오더등록/패킹','SCR-032 마스터오더상세'], RGBColor(0xE6,0x7E,0x22)),
    ('창고관리', ['SCR-040 입고처리','SCR-041 출고처리','SCR-042 창고현황'], RGBColor(0x8E,0x44,0xAD)),
    ('운송Tracking', ['SCR-050 AIR Tracking','SCR-051 SEA Tracking','SCR-052 CIR Tracking','SCR-053 통관신고현황'], RGBColor(0x16,0xA0,0x85)),
    ('회계/청구', ['SCR-060 청구서목록','SCR-061 청구서상세','SCR-062 입금처리','SCR-063 세금계산서','SCR-064 수입/비용현황','SCR-065 운송원가관리'], RGBColor(0xD3,0x54,0x00)),
    ('VOC관리', ['SCR-070 VOC목록','SCR-071 VOC등록','SCR-072 VOC답변'], RGBColor(0x2C,0x3E,0x50)),
    ('고객지원', ['SCR-080 QnA목록/등록','SCR-081 QnA답변','SCR-082 FAQ','SCR-083 공지사항'], RGBColor(0x7F,0x8C,0x8D)),
    ('관리자', ['SCR-090~099 (대시보드,회원,법인심사,메뉴,코드,권한,택배사,알림,통계,백업)'], RED),
]

col_w = 6.25
for i, (cat, items, c) in enumerate(summary):
    col = i % 2
    row = i // 2
    lx = 0.3 + col * col_w
    ty = 1.05 + row * 1.2
    if i == 9:  # 관리자 - 마지막
        lx = 0.3
        ty = 1.05 + 5 * 1.2
    add_rect(slide, lx, ty, col_w-0.15, 0.28, fill=c)
    add_text_box(slide, cat, lx+0.1, ty, col_w-0.25, 0.28, size=8, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text_box(slide, f'• {item}', lx+0.15, ty+0.3+j*0.2, col_w-0.3, 0.2, size=7.5, color=BLACK)

prs.save('SNTL_UI화면설계서.pptx')
print(f'PPT saved. Total slides: {len(prs.slides)}')
